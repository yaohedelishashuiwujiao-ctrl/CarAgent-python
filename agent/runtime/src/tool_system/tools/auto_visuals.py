from __future__ import annotations

import math
import os
import re
from pathlib import Path
from typing import Any

from ..context import ToolContext
from ..errors import ToolInputError
from ..protocol import ToolOutcomeStatus, ToolResult
from ..preflight import PreflightDecision
from ..registry import ToolCapability, ToolDependencies, ToolExecutionPolicy, ToolSpec
from .subjects_sql import MAX_SQL_LIMIT, _json_safe, execute_subjects_read_query, preflight_subjects_read_query


CHART_TYPES = {"scatter", "line", "bar", "grouped_bar", "pie", "donut", "histogram", "box", "heatmap"}
OUTPUT_DIR_NAME = "outputs"


def _slug(value: str, default: str) -> str:
    text = re.sub(r"[^\w.\-]+", "_", value.strip(), flags=re.UNICODE)
    text = text.strip("._")
    return text or default


def _output_path(context: ToolContext, file_name: str, suffix: str) -> Path:
    name = _slug(file_name, f"output{suffix}")
    if not name.lower().endswith(suffix):
        name = f"{name}{suffix}"
    base = Path(os.getenv("CLAWD_OUTPUT_DIR") or context.permission_context.workspace_root or context.workspace_root)
    path = (base / name if base.name == OUTPUT_DIR_NAME else base / OUTPUT_DIR_NAME / name).resolve()
    context.ensure_allowed_path(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def _file_url(path: Path, context: ToolContext) -> str | None:
    base = Path(os.getenv("CLAWD_OUTPUT_DIR") or context.permission_context.workspace_root or context.workspace_root)
    output_dir = (base if base.name == OUTPUT_DIR_NAME else base / OUTPUT_DIR_NAME).resolve()
    try:
        rel = path.resolve().relative_to(output_dir)
    except ValueError:
        return None
    return "/outputs/" + "/".join(rel.parts)


def _rows_from_input(
    tool_input: dict[str, Any],
    context: ToolContext,
    default_limit: int = MAX_SQL_LIMIT,
) -> tuple[list[dict[str, Any]], dict[str, Any], ToolResult | None]:
    if "data" in tool_input and tool_input["data"] is not None:
        data = tool_input["data"]
        if not isinstance(data, list):
            raise ToolInputError("data must be a list of objects")
        rows: list[dict[str, Any]] = []
        for item in data:
            if not isinstance(item, dict):
                raise ToolInputError("each data item must be an object")
            rows.append(item)
        return rows, {"source_type": "inline_data", "row_count": len(rows)}, None

    sql = str(tool_input.get("sql_query") or "").strip()
    if not sql:
        raise ToolInputError("provide either data or sql_query")
    limit = max(1, min(int(tool_input.get("limit") or default_limit), MAX_SQL_LIMIT))
    query_result = execute_subjects_read_query(
        sql,
        context,
        limit=limit,
        result_name="AutoChartGenerate",
    )
    if query_result.is_error:
        return [], {}, query_result
    output = query_result.output if isinstance(query_result.output, dict) else {}
    rows = output.get("rows") if isinstance(output.get("rows"), list) else []
    provenance = {
        "source_type": "subjects_sql",
        "query": output.get("query"),
        "tables": output.get("tables"),
        "estimated_rows": output.get("estimated_rows"),
        "evidence_hash": output.get("evidence_hash"),
        "row_count": output.get("row_count"),
    }
    return rows, provenance, None


def _numeric(value: Any) -> float | None:
    if value is None:
        return None
    if isinstance(value, (int, float)) and not isinstance(value, bool):
        if math.isfinite(float(value)):
            return float(value)
        return None
    text = str(value).replace(",", "").strip()
    match = re.search(r"-?\d+(?:\.\d+)?", text)
    if not match:
        return None
    try:
        return float(match.group(0))
    except ValueError:
        return None


def _clean_pairs(rows: list[dict[str, Any]], x: str, y: str) -> tuple[list[float], list[float], list[dict[str, Any]]]:
    xs: list[float] = []
    ys: list[float] = []
    used: list[dict[str, Any]] = []
    for row in rows:
        xv = _numeric(row.get(x))
        yv = _numeric(row.get(y))
        if xv is None or yv is None:
            continue
        xs.append(xv)
        ys.append(yv)
        used.append(row)
    return xs, ys, used


def _configure_matplotlib() -> Any:
    mpl_config = Path(os.getenv("MPLCONFIGDIR") or "/tmp/clawd-matplotlib")
    mpl_config.mkdir(parents=True, exist_ok=True)
    os.environ.setdefault("MPLCONFIGDIR", str(mpl_config))

    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib import font_manager

    cjk_font_paths: list[Path] = []
    for root in (Path("/usr/share/fonts"), Path("/usr/local/share/fonts"), Path.home() / ".local/share/fonts"):
        if not root.exists():
            continue
        for pattern in ("*NotoSansCJK*.ttc", "*NotoSansCJK*.otf", "*NotoSansCJK*.ttf", "*SourceHanSans*.otf", "*SourceHanSans*.ttf", "*wqy*.ttf", "*WenQuanYi*.ttf"):
            cjk_font_paths.extend(root.rglob(pattern))
    for font_path in cjk_font_paths[:20]:
        try:
            font_manager.fontManager.addfont(str(font_path))
        except Exception:
            pass

    # Force font cache rebuild after addfont; stale cache commonly breaks CJK rendering.
    try:
        font_manager._load_fontmanager(try_read_cache=False)
    except Exception:
        pass

    preferred = (
        "Noto Sans CJK SC",
        "Noto Sans CJK",
        "WenQuanYi Micro Hei",
        "Microsoft YaHei",
        "SimHei",
        "Arial Unicode MS",
        "DejaVu Sans",
    )
    names = {f.name for f in font_manager.fontManager.ttflist}
    chosen = None
    for name in preferred:
        if name in names:
            chosen = name
            break
    if chosen is None and cjk_font_paths:
        try:
            chosen = font_manager.FontProperties(fname=str(cjk_font_paths[0])).get_name()
        except Exception:
            chosen = None
    if chosen:
        plt.rcParams["font.family"] = "sans-serif"
        plt.rcParams["font.sans-serif"] = [chosen, "Noto Sans", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False
    try:
        plt.style.use("seaborn-v0_8-whitegrid")
    except Exception:
        pass
    return plt


def _apply_cjk_font(fig: Any) -> None:
    from matplotlib import font_manager

    preferred = ("Noto Sans CJK SC", "Noto Sans Mono CJK SC", "Noto Sans CJK", "WenQuanYi Micro Hei", "SimHei", "Microsoft YaHei")
    names = {f.name for f in font_manager.fontManager.ttflist}
    chosen = next((name for name in preferred if name in names), None)
    if not chosen:
        return
    prop = font_manager.FontProperties(family=chosen)
    for text in fig.findobj(match=lambda obj: hasattr(obj, "set_fontproperties")):
        try:
            text.set_fontproperties(prop)
        except Exception:
            pass


def _image_size(path: Path) -> tuple[int, int]:
    from PIL import Image

    with Image.open(path) as im:
        return im.size


def _fit_box(image_size: tuple[int, int], box_size: tuple[float, float]) -> tuple[float, float]:
    iw, ih = image_size
    bw, bh = box_size
    if iw <= 0 or ih <= 0:
        return bw, bh
    scale = min(bw / iw, bh / ih)
    return iw * scale, ih * scale


class AutoChartGenerateTool:
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="AutoChartGenerate",
            description=(
                "Generate deterministic automotive analysis charts as PNG files. "
                "Use this instead of Bash/Write for charts. Supports scatter, line, bar, grouped_bar, "
                "pie, donut, histogram, box, and heatmap. Input can be data rows or a read-only sql_query. "
                "Always returns row_count, plotted_count, file_path, and URL."
            ),
            input_schema={
                "type": "object",
                "additionalProperties": False,
                "properties": {
                    "chart_type": {"type": "string", "enum": sorted(CHART_TYPES)},
                    "title": {"type": "string"},
                    "data": {"type": "array", "items": {"type": "object"}},
                    "sql_query": {"type": "string"},
                    "limit": {"type": "integer"},
                    "x": {"type": "string"},
                    "y": {"type": "string"},
                    "category": {"type": "string"},
                    "value": {"type": "string"},
                    "series": {"type": "string"},
                    "label": {"type": "string"},
                    "color_by": {"type": "string"},
                    "x_label": {"type": "string"},
                    "y_label": {"type": "string"},
                    "file_name": {"type": "string"},
                    "width": {"type": "number"},
                    "height": {"type": "number"},
                    "dpi": {"type": "integer"},
                    "top_n": {"type": "integer"},
                },
                "required": ["chart_type", "title"],
            },
            is_read_only=False,
            strict=True,
            max_result_size_chars=20_000,
            capability=ToolCapability(
                namespace="artifact.chart",
                actions=("render",),
                entity_types=("tabular_data",),
                input_modes=("inline_rows", "governed_mysql_select"),
                output_modes=("png", "artifact", "evidence_manifest"),
                limitations=(
                    "Supports nine deterministic chart types and at most 500 governed SQL rows.",
                    "Does not perform statistical inference or decide whether a chart is analytically valid.",
                ),
                positive_examples=("Render a bar or scatter chart from an existing governed result.",),
                negative_examples=("Do not use as an alternate SQL executor or data-cleaning engine.",),
            ),
            execution=ToolExecutionPolicy(
                risk="low",
                side_effect="artifact",
                timeout_s=30,
                concurrency_pool="artifact",
                supports_parallel=True,
                idempotent=False,
            ),
            dependencies=ToolDependencies(
                services=("artifact_storage", "subjects_mysql_if_sql_input"),
                health_probe="chart_dependencies_health",
            ),
            preflight_checks=(
                "tool_authorized",
                "artifact_workspace",
                "chart_input_shape",
                "governed_sql_if_present",
            ),
        )

    def run(self, tool_input: dict[str, Any], context: ToolContext) -> ToolResult:
        chart_type = str(tool_input["chart_type"])
        rows, provenance, query_error = _rows_from_input(tool_input, context)
        if query_error is not None:
            return ToolResult(
                name="AutoChartGenerate",
                output={
                    "error": "chart data query was rejected by the governed SQL executor",
                    "query_error": query_error.output,
                },
                is_error=True,
                outcome_status=query_error.outcome_status,
                reason_code=query_error.reason_code,
                retryable=query_error.retryable,
                diagnostics=query_error.diagnostics,
            )
        if not rows:
            return ToolResult(
                name="AutoChartGenerate",
                output={"error": "no rows to plot", "data_source": provenance},
                is_error=True,
                outcome_status=ToolOutcomeStatus.NO_DATA,
                reason_code="CHART_NO_ROWS",
            )

        title = str(tool_input.get("title") or chart_type)
        file_name = str(tool_input.get("file_name") or _slug(title, chart_type))
        path = _output_path(context, file_name, ".png")
        width = float(tool_input.get("width") or 11)
        height = float(tool_input.get("height") or 6.2)
        dpi = max(72, min(int(tool_input.get("dpi") or 160), 300))
        top_n = max(1, min(int(tool_input.get("top_n") or 30), 80))

        try:
            plt = _configure_matplotlib()
        except Exception as exc:
            return ToolResult(name="AutoChartGenerate", output={"error": f"matplotlib unavailable: {exc}"}, is_error=True)

        fig, ax = plt.subplots(figsize=(width, height), dpi=dpi)
        plotted_count = 0
        warnings: list[str] = []

        def require(name: str) -> str:
            value = str(tool_input.get(name) or "").strip()
            if not value:
                raise ToolInputError(f"{chart_type} requires '{name}'")
            return value

        try:
            if chart_type == "scatter":
                x = require("x")
                y = require("y")
                xs, ys, used = _clean_pairs(rows, x, y)
                if not xs:
                    raise ToolInputError(f"no numeric points for x={x}, y={y}")
                color_by = str(tool_input.get("color_by") or "").strip()
                if color_by:
                    categories = [str(row.get(color_by) or "未知") for row in used]
                    unique = list(dict.fromkeys(categories))
                    cmap = plt.get_cmap("tab20")
                    color_map = {cat: cmap(i % 20) for i, cat in enumerate(unique)}
                    ax.scatter(xs, ys, c=[color_map[c] for c in categories], alpha=0.72, s=26, edgecolors="none")
                    if len(unique) <= 12:
                        handles = [
                            plt.Line2D([0], [0], marker="o", color="w", markerfacecolor=color_map[c], label=c, markersize=7)
                            for c in unique
                        ]
                        ax.legend(handles=handles, fontsize=8, loc="best")
                else:
                    ax.scatter(xs, ys, alpha=0.72, s=26, color="#2563eb", edgecolors="none")
                plotted_count = len(xs)
                ax.set_xlabel(str(tool_input.get("x_label") or x))
                ax.set_ylabel(str(tool_input.get("y_label") or y))
                ax.grid(True, alpha=0.25)

            elif chart_type in {"line", "bar"}:
                x = require("x")
                y = require("y")
                pairs = [(str(row.get(x) or ""), _numeric(row.get(y))) for row in rows]
                pairs = [(a, b) for a, b in pairs if a and b is not None]
                # Sort by value descending for bar charts, keep original order for line
                if chart_type == "bar":
                    pairs.sort(key=lambda p: p[1], reverse=True)
                pairs = pairs[:top_n]
                if not pairs:
                    raise ToolInputError(f"no values for x={x}, y={y}")
                labels = [a for a, _ in pairs]
                values = [float(b) for _, b in pairs]
                if chart_type == "line":
                    ax.plot(labels, values, marker="o", color="#0f766e", linewidth=2)
                else:
                    ax.bar(labels, values, color="#2563eb")
                plotted_count = len(values)
                ax.set_xlabel(str(tool_input.get("x_label") or x))
                ax.set_ylabel(str(tool_input.get("y_label") or y))
                ax.tick_params(axis="x", rotation=35, labelsize=8)
                ax.grid(True, axis="y", alpha=0.25)

            elif chart_type in {"pie", "donut"}:
                category = str(tool_input.get("category") or tool_input.get("label") or "").strip()
                value = require("value")
                if not category:
                    raise ToolInputError("pie/donut requires 'category' or 'label'")
                items = [(str(row.get(category) or "未知"), _numeric(row.get(value))) for row in rows]
                items = [(a, b) for a, b in items if b is not None and b > 0]
                # Sort by value descending before truncation
                items.sort(key=lambda p: p[1], reverse=True)
                items = items[:top_n]
                if not items:
                    raise ToolInputError(f"no positive values for value={value}")
                labels = [a for a, _ in items]
                values = [float(b) for _, b in items]
                wedgeprops = {"width": 0.42} if chart_type == "donut" else None
                ax.pie(values, labels=labels, autopct="%1.1f%%", startangle=90, wedgeprops=wedgeprops, textprops={"fontsize": 8})
                ax.axis("equal")
                plotted_count = len(values)

            elif chart_type == "histogram":
                value = require("value")
                # Detect pre-aggregated data (e.g., [{"range": "<2800", "cnt": 13}, ...])
                # If rows are few (<50) and a count-like column exists, this is likely pre-binned
                count_cols = {"cnt", "count", "frequency", "num", "total", "vehicle_count"}
                has_count_col = any(c in count_cols for c in rows[0].keys()) if rows else False
                if has_count_col and len(rows) < 50:
                    # Pre-aggregated: use bar chart instead of re-binning
                    warnings.append(f"data appears pre-aggregated ({len(rows)} buckets); rendering as bar chart")
                    cnt_col = next((c for c in count_cols if c in rows[0].keys()), "cnt")
                    label_col = next((c for c in rows[0].keys() if c != cnt_col), "label")
                    pairs = [(str(row.get(label_col) or str(i)), _numeric(row.get(cnt_col))) for i, row in enumerate(rows)]
                    pairs = [(a, b) for a, b in pairs if a and b is not None]
                    pairs.sort(key=lambda p: p[1], reverse=True)
                    pairs = pairs[:top_n]
                    if not pairs:
                        raise ToolInputError(f"no numeric values for {cnt_col}")
                    labels = [a for a, _ in pairs]
                    vals = [float(b) for _, b in pairs]
                    ax.bar(labels, vals, color="#0f766e")
                    ax.set_xlabel(str(tool_input.get("x_label") or label_col))
                    ax.set_ylabel(str(tool_input.get("y_label") or "数量"))
                    ax.tick_params(axis="x", rotation=35, labelsize=7)
                    ax.grid(True, axis="y", alpha=0.25)
                    plotted_count = len(vals)
                else:
                    values = [_numeric(row.get(value)) for row in rows]
                    values = [v for v in values if v is not None]
                    if not values:
                        raise ToolInputError(f"no numeric values for value={value}")
                    if len(values) < 3:
                        warnings.append(f"very few numeric values ({len(values)}); histogram may be misleading")
                    ax.hist(values, bins=min(30, max(8, int(math.sqrt(len(values))))), color="#0f766e", alpha=0.82)
                    ax.set_xlabel(str(tool_input.get("x_label") or value))
                    ax.set_ylabel(str(tool_input.get("y_label") or "数量"))
                    ax.grid(True, axis="y", alpha=0.25)
                    plotted_count = len(values)

            elif chart_type == "box":
                value = str(tool_input.get("value") or tool_input.get("y") or tool_input.get("x") or "").strip()
                if not value:
                    raise ToolInputError("box requires 'value' or 'y'")
                category = str(tool_input.get("category") or "").strip()
                if category:
                    grouped: dict[str, list[float]] = {}
                    for row in rows:
                        val = _numeric(row.get(value))
                        if val is None:
                            continue
                        grouped.setdefault(str(row.get(category) or "未知"), []).append(val)
                    grouped = {k: v for k, v in grouped.items() if v}
                    items = list(grouped.items())[: min(top_n, 20)]
                    if not items:
                        raise ToolInputError("no grouped numeric values")
                    ax.boxplot([v for _, v in items], tick_labels=[k for k, _ in items], showfliers=True)
                    ax.tick_params(axis="x", rotation=35, labelsize=8)
                    plotted_count = sum(len(v) for _, v in items)
                else:
                    values = [_numeric(row.get(value)) for row in rows]
                    values = [v for v in values if v is not None]
                    if not values:
                        raise ToolInputError(f"no numeric values for value={value}")
                    ax.boxplot(values, tick_labels=[value], showfliers=True)
                    plotted_count = len(values)
                ax.set_ylabel(str(tool_input.get("y_label") or value))
                ax.grid(True, axis="y", alpha=0.25)

            elif chart_type in {"grouped_bar", "heatmap"}:
                x = require("x")
                y = require("y")
                value = require("value")
                import pandas as pd

                df = pd.DataFrame(rows)
                if x not in df.columns or y not in df.columns or value not in df.columns:
                    raise ToolInputError("grouped_bar/heatmap fields must exist in rows")
                df[value] = pd.to_numeric(df[value], errors="coerce")
                pivot = df.pivot_table(index=x, columns=y, values=value, aggfunc="mean")
                # Sort by row sum descending, then take top_n
                pivot["_sort_key"] = pivot.sum(axis=1)
                pivot = pivot.sort_values("_sort_key", ascending=False).head(top_n)
                pivot = pivot.drop(columns=["_sort_key"])
                if pivot.empty:
                    raise ToolInputError("pivot table is empty")
                if chart_type == "grouped_bar":
                    pivot.plot(kind="bar", ax=ax)
                    ax.tick_params(axis="x", rotation=35, labelsize=8)
                    ax.legend(fontsize=8)
                else:
                    im = ax.imshow(pivot.values, aspect="auto", cmap="YlGnBu")
                    ax.set_xticks(range(len(pivot.columns)), [str(c) for c in pivot.columns], rotation=35, ha="right", fontsize=8)
                    ax.set_yticks(range(len(pivot.index)), [str(i) for i in pivot.index], fontsize=8)
                    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
                plotted_count = int(pivot.count().sum())

            else:
                raise ToolInputError(f"unsupported chart_type: {chart_type}")

            ax.set_title(f"{title}  (n={plotted_count})", fontsize=14, pad=12)
            _apply_cjk_font(fig)
            fig.tight_layout()
            fig.savefig(path)
            plt.close(fig)
        except Exception:
            plt.close(fig)
            raise

        return ToolResult(
            name="AutoChartGenerate",
            output=_json_safe(
                {
                    "file_path": str(path),
                    "url": _file_url(path, context),
                    "chart_type": chart_type,
                    "row_count": len(rows),
                    "plotted_count": plotted_count,
                    "warnings": warnings,
                    "data_source": provenance,
                }
            ),
        )

    def preflight(self, tool_input: dict[str, Any], context: ToolContext) -> PreflightDecision:
        if tool_input.get("data") is not None:
            data = tool_input.get("data")
            if not isinstance(data, list) or not data:
                return PreflightDecision.reject("CHART_NO_ROWS", "Chart data must be a non-empty array.")
            return PreflightDecision.allow("CHART_INLINE_DATA_READY")
        sql = str(tool_input.get("sql_query") or "").strip()
        if not sql:
            return PreflightDecision.reject(
                "CHART_DATA_SOURCE_MISSING",
                "Provide either non-empty data rows or a governed sql_query.",
            )
        return preflight_subjects_read_query(sql, context)


class AutoPptxGenerateTool:
    def spec(self) -> ToolSpec:
        return ToolSpec(
            name="AutoPptxGenerate",
            description=(
                "Generate a deterministic PowerPoint deck. requested_slide_count is mandatory and must exactly "
                "match len(slides). Never add pages automatically. Each slide should state one conclusion and can "
                "include one chart/image, a small table, key_points, and a bottom conclusion."
            ),
            input_schema={
                "type": "object",
                "additionalProperties": False,
                "properties": {
                    "requested_slide_count": {"type": "integer", "minimum": 1, "maximum": 20},
                    "deck_title": {"type": "string", "minLength": 1},
                    "file_name": {"type": "string"},
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "additionalProperties": False,
                            "properties": {
                                "title": {"type": "string", "minLength": 1},
                                "subtitle": {"type": "string"},
                                "chart_path": {"type": "string"},
                                "image_path": {"type": "string"},
                                "key_points": {
                                    "type": "array",
                                    "items": {"type": "string"},
                                    "maxItems": 5,
                                    "description": "At most five concise findings; put detailed comparisons in table or notes.",
                                },
                                "conclusion": {"type": "string", "minLength": 1},
                                "table": {
                                    "type": "object",
                                    "additionalProperties": False,
                                    "properties": {
                                        "columns": {"type": "array", "items": {"type": "string"}},
                                        "rows": {
                                            "type": "array",
                                            "items": {"type": "array"},
                                            "maxItems": 8,
                                        },
                                    },
                                },
                                "notes": {"type": "string"},
                                "source_footer": {
                                    "type": "string",
                                    "description": "Short visible source names or URLs rendered in the slide footer.",
                                },
                            },
                            "required": ["title", "conclusion"],
                        },
                    },
                },
                "required": ["requested_slide_count", "deck_title", "slides"],
            },
            is_read_only=False,
            strict=True,
            max_result_size_chars=20_000,
            capability=ToolCapability(
                namespace="artifact.pptx",
                actions=("render", "package"),
                entity_types=("analysis_deck",),
                input_modes=("slide_spec", "chart_artifact", "image_artifact"),
                output_modes=("pptx", "artifact"),
                limitations=(
                    "Supports 1-20 explicit slides using a deterministic layout.",
                    "Does not independently verify factual claims or citation coverage.",
                ),
                positive_examples=("Generate an exact six-slide deck after evidence and charts are ready.",),
                negative_examples=("Do not call before required evidence and media artifacts exist.",),
            ),
            execution=ToolExecutionPolicy(
                risk="low",
                side_effect="artifact",
                timeout_s=60,
                concurrency_pool="artifact",
                supports_parallel=False,
                idempotent=False,
            ),
            dependencies=ToolDependencies(
                services=("artifact_storage",),
                health_probe="pptx_dependencies_health",
            ),
            preflight_checks=(
                "tool_authorized",
                "artifact_workspace",
                "slide_count_contract",
                "media_artifacts_exist",
            ),
        )

    def run(self, tool_input: dict[str, Any], context: ToolContext) -> ToolResult:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except Exception as exc:
            return ToolResult(name="AutoPptxGenerate", output={"error": f"python-pptx unavailable: {exc}"}, is_error=True)

        requested = int(tool_input["requested_slide_count"])
        slides = tool_input.get("slides") or []
        if not isinstance(slides, list):
            raise ToolInputError("slides must be an array")
        if requested < 1 or requested > 20:
            raise ToolInputError("requested_slide_count must be between 1 and 20")
        if len(slides) != requested:
            return ToolResult(
                name="AutoPptxGenerate",
                output={
                    "error": "slide count mismatch",
                    "requested_slide_count": requested,
                    "provided_slide_count": len(slides),
                    "rule": "The tool never adds or removes slides automatically. Compress or expand the slides array explicitly.",
                },
                is_error=True,
            )

        for idx, slide in enumerate(slides, start=1):
            key_points = slide.get("key_points") or []
            table = slide.get("table") or {}
            rows = table.get("rows") or [] if isinstance(table, dict) else []
            if len(key_points) > 5:
                return ToolResult(name="AutoPptxGenerate", output={"error": f"slide {idx} has too many key_points; max 5"}, is_error=True)
            if len(rows) > 8:
                return ToolResult(name="AutoPptxGenerate", output={"error": f"slide {idx} table has too many rows; max 8"}, is_error=True)

        deck_title = str(tool_input.get("deck_title") or "Analysis")
        file_name = str(tool_input.get("file_name") or _slug(deck_title, "deck"))
        path = _output_path(context, file_name, ".pptx")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        navy = RGBColor(31, 41, 55)
        teal = RGBColor(15, 118, 110)
        muted = RGBColor(100, 116, 139)
        line = RGBColor(226, 232, 240)
        white = RGBColor(255, 255, 255)
        panel = RGBColor(248, 250, 252)

        for slide_data in slides:
            slide = prs.slides.add_slide(blank)
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = panel
            bg.line.fill.background()

            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.14))
            accent.fill.solid()
            accent.fill.fore_color.rgb = teal
            accent.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.0), Inches(0.52))
            title_tf = title_box.text_frame
            title_tf.clear()
            p = title_tf.paragraphs[0]
            p.text = str(slide_data["title"])
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = navy

            subtitle = str(slide_data.get("subtitle") or "").strip()
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.82), Inches(12.0), Inches(0.34))
                sub = sub_box.text_frame.paragraphs[0]
                sub.text = subtitle
                sub.font.name = "Microsoft YaHei"
                sub.font.size = Pt(10)
                sub.font.color.rgb = muted

            media_path = str(slide_data.get("chart_path") or slide_data.get("image_path") or "").strip()
            has_media = False
            if media_path:
                media = context.ensure_allowed_path(media_path)
                if not media.exists():
                    return ToolResult(name="AutoPptxGenerate", output={"error": f"image not found: {media}"}, is_error=True)
                media_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.18), Inches(7.55), Inches(4.82))
                media_card.fill.solid()
                media_card.fill.fore_color.rgb = white
                media_card.line.color.rgb = line
                tag = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(1.2), Inches(0.25))
                tag_tf = tag.text_frame
                tag_tf.clear()
                tag_p = tag_tf.paragraphs[0]
                tag_p.text = "图表"
                tag_p.font.name = "Microsoft YaHei"
                tag_p.font.size = Pt(9)
                tag_p.font.bold = True
                tag_p.font.color.rgb = teal
                img_w, img_h = _image_size(media)
                fit_w, fit_h = _fit_box((img_w, img_h), (6.95 * 96, 4.0 * 96))
                slide.shapes.add_picture(
                    str(media),
                    Inches(0.78 + (6.95 - fit_w / 96) / 2),
                    Inches(1.6 + (4.0 - fit_h / 96) / 2),
                    width=Inches(fit_w / 96),
                    height=Inches(fit_h / 96),
                )
                has_media = True

            left = Inches(8.35) if has_media else Inches(0.55)
            width = Inches(4.3) if has_media else Inches(12.2)
            y = Inches(1.18)

            info_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                Inches(1.18),
                width,
                Inches(4.82),
            )
            info_card.fill.solid()
            info_card.fill.fore_color.rgb = white
            info_card.line.color.rgb = line

            inner_left = left + Inches(0.18)
            inner_width = width - Inches(0.36)
            y = Inches(1.4)

            key_points = [str(x) for x in (slide_data.get("key_points") or [])]
            if key_points:
                head = slide.shapes.add_textbox(inner_left, y, inner_width, Inches(0.25))
                head_tf = head.text_frame
                head_tf.clear()
                head_p = head_tf.paragraphs[0]
                head_p.text = "关键发现"
                head_p.font.name = "Microsoft YaHei"
                head_p.font.size = Pt(10)
                head_p.font.bold = True
                head_p.font.color.rgb = teal
                box = slide.shapes.add_textbox(inner_left, y + Inches(0.28), inner_width, Inches(1.55))
                tf = box.text_frame
                tf.clear()
                for i, item in enumerate(key_points):
                    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    para.text = f"• {item}"
                    para.level = 0
                    para.font.name = "Microsoft YaHei"
                    para.font.size = Pt(13)
                    para.font.color.rgb = navy
                    para.space_after = Pt(4)
                y = Inches(3.12)

            table = slide_data.get("table") or {}
            if isinstance(table, dict) and table.get("columns") and table.get("rows"):
                columns = [str(c) for c in table["columns"]]
                rows = table["rows"]
                head = slide.shapes.add_textbox(inner_left, y, inner_width, Inches(0.25))
                head_tf = head.text_frame
                head_tf.clear()
                head_p = head_tf.paragraphs[0]
                head_p.text = "数据摘要"
                head_p.font.name = "Microsoft YaHei"
                head_p.font.size = Pt(10)
                head_p.font.bold = True
                head_p.font.color.rgb = teal
                table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), inner_left, y + Inches(0.28), inner_width, Inches(1.55)).table
                for col_idx, col in enumerate(columns):
                    cell = table_shape.cell(0, col_idx)
                    cell.text = col
                    p0 = cell.text_frame.paragraphs[0]
                    p0.font.name = "Microsoft YaHei"
                    p0.font.bold = True
                    p0.font.size = Pt(9)
                    p0.font.color.rgb = white
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = teal
                for row_idx, row in enumerate(rows, start=1):
                    for col_idx, value in enumerate(list(row)[: len(columns)]):
                        cell = table_shape.cell(row_idx, col_idx)
                        cell.text = str(value)
                        p0 = cell.text_frame.paragraphs[0]
                        p0.font.name = "Microsoft YaHei"
                        p0.font.size = Pt(8)
                        p0.font.color.rgb = navy
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = white if row_idx % 2 else RGBColor(241, 245, 249)

            concl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(6.12), Inches(12.18), Inches(0.72))
            concl.fill.solid()
            concl.fill.fore_color.rgb = teal
            concl.line.fill.background()
            tf = concl.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "结论：" + str(slide_data["conclusion"])
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = white
            p.alignment = PP_ALIGN.CENTER

            source_footer = str(slide_data.get("source_footer") or "").strip()
            if source_footer:
                source_box = slide.shapes.add_textbox(
                    Inches(0.62), Inches(6.93), Inches(12.0), Inches(0.28)
                )
                source_tf = source_box.text_frame
                source_tf.clear()
                source_p = source_tf.paragraphs[0]
                source_p.text = "来源：" + source_footer
                source_p.font.name = "Microsoft YaHei"
                source_p.font.size = Pt(7)
                source_p.font.color.rgb = muted

            notes = str(slide_data.get("notes") or "").strip()
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        prs.save(path)
        return ToolResult(
            name="AutoPptxGenerate",
            output=_json_safe(
                {
                    "file_path": str(path),
                    "url": _file_url(path, context),
                    "slide_count": len(slides),
                    "requested_slide_count": requested,
                    "deck_title": deck_title,
                    "slide_manifest": [
                        {
                            "title": str(slide.get("title") or ""),
                            "has_table": bool(
                                isinstance(slide.get("table"), dict)
                                and (slide.get("table") or {}).get("rows")
                            ),
                            "has_media": bool(slide.get("chart_path") or slide.get("image_path")),
                            "has_source_footer": bool(str(slide.get("source_footer") or "").strip()),
                        }
                        for slide in slides
                    ],
                }
            ),
        )

    def preflight(self, tool_input: dict[str, Any], context: ToolContext) -> PreflightDecision:
        requested = int(tool_input.get("requested_slide_count") or 0)
        slides = tool_input.get("slides")
        output_contract = context.runtime_state.get("output_contract")
        contract_slide_count = None
        for artifact in getattr(output_contract, "required_artifacts", ()):
            if getattr(artifact, "artifact_type", None) == "pptx":
                contract_slide_count = getattr(artifact, "exact_count", None)
                break
        if contract_slide_count is not None and requested != contract_slide_count:
            return PreflightDecision.reject(
                "PPTX_OUTPUT_CONTRACT_MISMATCH",
                f"The user requires exactly {contract_slide_count} slides; requested_slide_count was {requested}.",
                diagnostics={"contract_slide_count": contract_slide_count, "requested_slide_count": requested},
            )
        if requested < 1 or requested > 20:
            return PreflightDecision.reject(
                "PPTX_SLIDE_COUNT_OUT_OF_RANGE",
                "requested_slide_count must be between 1 and 20.",
            )
        if not isinstance(slides, list) or len(slides) != requested:
            return PreflightDecision.reject(
                "PPTX_SLIDE_COUNT_MISMATCH",
                f"Expected exactly {requested} slides, received {len(slides) if isinstance(slides, list) else 0}.",
                diagnostics={"requested": requested, "provided": len(slides) if isinstance(slides, list) else None},
            )
        for index, slide in enumerate(slides, start=1):
            if not isinstance(slide, dict):
                return PreflightDecision.reject("PPTX_SLIDE_INVALID", f"Slide {index} must be an object.")
            key_points = slide.get("key_points") or []
            table = slide.get("table") or {}
            rows = table.get("rows") or [] if isinstance(table, dict) else []
            if not isinstance(key_points, list) or len(key_points) > 5:
                return PreflightDecision.reject(
                    "PPTX_KEY_POINTS_LIMIT",
                    f"Slide {index} key_points must contain at most 5 items.",
                    diagnostics={"slide": index, "maximum": 5},
                )
            if not isinstance(rows, list) or len(rows) > 8:
                return PreflightDecision.reject(
                    "PPTX_TABLE_ROWS_LIMIT",
                    f"Slide {index} table rows must contain at most 8 items.",
                    diagnostics={"slide": index, "maximum": 8},
                )
            requirement = next(
                (
                    item
                    for item in getattr(context.runtime_state.get("output_contract"), "required_artifacts", ())
                    if getattr(item, "artifact_type", None) == "pptx"
                ),
                None,
            )
            if requirement and getattr(requirement, "require_table_per_slide", False) and not rows:
                return PreflightDecision.reject(
                    "PPTX_TABLE_REQUIRED_PER_SLIDE",
                    f"Slide {index} requires a non-empty comparison/data table.",
                    diagnostics={"slide": index},
                )
            if requirement and getattr(requirement, "require_source_per_slide", False) and not str(slide.get("source_footer") or "").strip():
                return PreflightDecision.reject(
                    "PPTX_SOURCE_FOOTER_REQUIRED",
                    f"Slide {index} requires a visible source_footer.",
                    diagnostics={"slide": index},
                )
            media_path = str(slide.get("chart_path") or slide.get("image_path") or "").strip()
            if not media_path:
                continue
            try:
                media = context.ensure_allowed_path(media_path)
            except Exception as exc:
                return PreflightDecision.reject(
                    "PPTX_MEDIA_PATH_DENIED",
                    f"Slide {index} media path is not allowed: {exc}",
                )
            if not media.exists() or not media.is_file():
                return PreflightDecision.reject(
                    "PPTX_MEDIA_NOT_FOUND",
                    f"Slide {index} media artifact does not exist: {media}",
                )
        return PreflightDecision.allow("PPTX_INPUT_CONTRACT_PASSED")
